# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# import os




# def ppt_report_generation(image_folder: str, output_path: str = None):

#     prs = Presentation()
#     # 1. First slide — Title Slide layout
#     title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout 0: Title Slide
#     title = title_slide.shapes.title
#     subtitle = title_slide.placeholders[1]

#     title.text = "Bounding Box Report"  # Header
#     subtitle.text = "File Name: File Name\nDescription: Description"
#     # Optionally style it:
#     for paragraph in subtitle.text_frame.paragraphs:
#         paragraph.font.size = Pt(18)
#     # — Source examples show setting title/subtitle via placeholders :contentReference[oaicite:0]{index=0}

#     # 2. Add image slides as before
#     blank_layout = prs.slide_layouts[6]
#     image_files = sorted(
#         f for f in os.listdir(image_folder)
#         if os.path.isfile(os.path.join(image_folder, f))
#            and f.lower().endswith(('.png', '.jpg', '.jpeg'))
#     )

#     for fname in image_files:
#         name, ext = os.path.splitext(fname)
#         slide = prs.slides.add_slide(blank_layout)
        
#         # Add title showing filename and format
#         if slide.shapes.title:
#             slide.shapes.title.text = f"{name} — {ext.lstrip('.').upper()}"
#         else:
#             tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
#             tf = tx_box.text_frame
#             p = tf.paragraphs[0]
#             p.text = f"{name} — {ext.lstrip('.').upper()}"
#             p.font.size = Pt(24)
#             p.font.bold = True
#             p.font.color.rgb = RGBColor(0, 0, 0)

#         # Insert image
#         img_path = os.path.join(image_folder, fname)
#         slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

#     prs.save(output_path)
#     # return path to saved powerpoint
#     return output_path



# ppt_report_generation(image_folder= r"C:\Users\kunakuma\Documents\AI CATIA Automation\AI-CATIA-Automation\static\catia_screenshots\BOUNDING_BOX_PROJECTED_AREA_POWER_COPY.CATPart", output_path= r"C:\Users\kunakuma\Downloads")
# # # image_folder = r"C:\Users\kunakuma\Documents\AI CATIA Automation\AI-CATIA-Automation\static\catia_screenshots\BOUNDING_BOX_PROJECTED_AREA_POWER_COPY.CATPart"
# # # output_pptx = r"C:\Users\kunakuma\Documents\AI CATIA Automation\image_report.pptx"

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from pathlib import Path

def ppt_report_generation(image_folder: str, filename: str) -> str:
    print("Generating PowerPoint report...")
    if not Path(image_folder).is_dir():
        raise FileNotFoundError(f"Invalid image folder: {image_folder}")

    # Define output path (Downloads folder)
    output_dir = Path.home() / "Downloads"
    output_path = output_dir / f"{filename}.pptx"
    print("Output PowerPoint path:", output_path)

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Bounding Box Report"
    subtitle = slide.placeholders[1]
    subtitle.text = f"File Name: {filename}\nDescription: Auto-generated report"
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(18)

    # Image slides
    blank_layout = prs.slide_layouts[6]
    image_files = sorted(f for f in os.listdir(image_folder)
                         if f.lower().endswith(('.png', '.jpg', '.jpeg')))

    for fname in image_files:
        slide = prs.slides.add_slide(blank_layout)
        title_text = f"{os.path.splitext(fname)[0]} — {os.path.splitext(fname)[1].lstrip('.').upper()}"

        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text_frame
        p = tx.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        img_path = os.path.join(image_folder, fname)
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

    prs.save(output_path)
    print(f"PPT saved at: {output_path}")
    return str(output_path)

#  — Usage example —
# image_folder = r"C:\Users\kunakuma\Documents\AI CATIA Automation\AI-CATIA-Automation\static\catia_screenshots\BOUNDING_BOX_PROJECTED_AREA_POWER_COPY.CATPart"
# output_file = str(Path.home() / "Downloads" / "image_report.pptx")

# try:
#     saved_path = ppt_report_generation(image_folder, output_file)
#     print("PPT saved at:", saved_path)
# except FileNotFoundError as e:
#     print("Error:", e)
